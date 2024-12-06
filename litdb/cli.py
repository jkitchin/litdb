"""CLI for litdb.

The main command is litdb. There are subcommands for the actions.
"""

import click
from rich import print as richprint
import os
from sentence_transformers import SentenceTransformer
import bs4
import requests
import datetime
import pathlib
import warnings
import tempfile
import sys

from tqdm import tqdm
import dateparser
import json
import numpy as np
import ollama
import time
import webbrowser
from docx import Document
from pptx import Presentation
import nbformat
from nbconvert import MarkdownExporter
from jinja2 import Template
import tabulate
from more_itertools import batched

from . import root, config
from .db import (get_db, add_source, add_work, add_author,
                 update_filter, add_bibtex)
from .openalex import get_data, get_text
from .pdf import add_pdf
from .bibtex import dump_bibtex


warnings.filterwarnings("ignore")

db = get_db()


@click.group()
def cli():
    """Group command for litdb."""
    pass


@cli.command()
def init():
    """Initialize a directory."""
    print(f'Initialized in {os.getcwd()}')

#################
# Add functions #
#################


@cli.command()
@click.argument('sources', nargs=-1)
@click.option('--references', is_flag=True, help='Add references too.')
@click.option('--related', is_flag=True, help='Add related too.')
@click.option('--citing', is_flag=True, help='Add citing too.')
@click.option('--all', is_flag=True, help='Add references, related and citing.')
def add(sources, references=False, citing=False, related=False, all=False):
    """Add WIDS to the db.

    SOURCES can be one or more of a doi or orcid, a pdf path, a url, bibtex
    file, or other kind of file assumed to be text.

    These are one time additions.

    """
    for source in tqdm(sources):

        # a work
        if source.startswith('10.') or 'doi.org' in source:
            if source.startswith('10.'):
                source = f'https://doi.org/{source}'

            if all:
                references, citing, related = True, True, True
            add_work(source, references, citing, related)

        # works from an author
        elif 'orcid' in source or source.lower().startswith('https://openalex.org/a'):
            add_author(source)

        # a bibtex file
        elif source.endswith('.bib'):
            add_bibtex(source)

        # pdf
        elif source.endswith('.pdf'):
            add_pdf(source)

        # docx
        elif source.endswith('.docx'):
            doc = Document(source)
            add_source(source, '\n'.join([para.text for para in doc.paragraphs]))

        # pptx
        elif source.endswith('.pptx'):
            prs = Presentation(source)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            add_source(source, '\n'.join(text))

        # html
        elif source.endswith('.html'):
            with open(source) as f:
                text = f.read()
            soup = bs4.BeautifulSoup(text, features="lxml")
            add_source(source, soup.get_text())

        # a url
        elif source.startswith('http'):
            soup = bs4.BeautifulSoup(requests.get(source).text)
            add_source(source, soup.get_text())

        # ipynb
        elif source.endswith('.ipynb'):
            with open(source, 'r', encoding='utf-8') as f:
                notebook = nbformat.read(f, as_version=4)

            # Create a Markdown exporter
            markdown_exporter = MarkdownExporter()

            # Convert the notebook to Markdown
            (body, resources) = markdown_exporter.from_notebook_node(notebook)

            add_source(source, body)

        # assume it is text
        else:
            with open(source) as f:
                text = f.read()
            add_source(source, text)


@cli.command()
@click.argument('sources', nargs=-1)
def index(sources):
    """Index the directories in SOURCES.

    SOURCES is a list of directories.
    """
    for directory in sources:
        directory = pathlib.Path(directory).resolve()
        for fname in directory.rglob('*'):
            # for f in files:
            if fname.suffix in ['.pdf', '.docx', '.pptx', '.org', '.md',
                                '.html', '.bib', '.ipynb']:
                fname = str(fname)

                # skip files we already have
                if db.execute('''select source from sources where source = ?''',
                              (fname,)).fetchone():
                    continue

                try:
                    print(f'Adding {fname}')
                    add([fname])
                    # add seems to call a SystemExit, so this line doesn't get
                    # run. I don't know why this happens here.
                    richprint(f'Added {fname}')
                # I don't know why this gets called, but it does, and I catch it
                # so we keep going
                except SystemExit:
                    pass
                except:
                    richprint('Still something worng')
                    import sys
                    richprint(sys.exc_info())

        last_updated = datetime.date.today().strftime('%Y-%m-%d')

        directory = str(directory)  # we need strings for the db
        if db.execute('''select path from directories where path = ?''',
                      (directory,)).fetchone():
            print(f'Updating {directory}')
            db.execute('''update directories set last_updated = ?
            where path = ?''', (last_updated, directory))
        else:
            print(f'Inserting {directory}: {last_updated}')
            db.execute('''insert into directories(path, last_updated)
            values (?, ?)''', (directory, last_updated))

        db.commit()


@cli.command()
def reindex():
    """Reindex saved directories."""
    for directory, in db.execute('''select path from directories''').fetchall():
        print(f'Reindexing {directory}')
        index([directory])


###########
# Tagging #
###########
@cli.command()
@click.argument('sources', nargs=-1)
@click.option('-t', '--tag', 'tags', multiple=True)
def add_tag(sources, tags):
    """Add tags to sources.

    It is a little annoying to add multiple tags. It looks like this.
    litdb add-tag source -t tag1 -t tag2
    """
    for source in sources:
        # Get source id
        source_id, = db.execute('select rowid from sources where source = ?',
                                (source,)).fetchone()

        for tag in tags:
            # get tag id
            tag_id = db.execute('select rowid from tags where tag = ?',
                                (tag,)).fetchone()

            if not tag_id:
                c = db.execute('insert into tags(tag) values (?)', (tag,))
                tag_id = c.lastrowid
                db.commit()
            else:
                # we get a tuple in the first query
                tag_id, = tag_id

            # Now add a tag
            db.execute('insert into source_tag(source_id, tag_id) values (?, ?)',
                       (source_id, tag_id))
            db.commit()

            print(f'Tagged {source} with {tag}')


@cli.command()
@click.argument('sources', nargs=-1)
@click.option('-t', '--tag', 'tags', multiple=True)
def rm_tag(sources, tags):
    """Remove tags from sources.

    It is a little annoying to remove multiple tags. It looks like this.
    litdb rm-tag source -t tag1 -t tag2
    """
    for source in sources:
        # Get source id
        source_id, = db.execute('select rowid from sources where source = ?',
                                (source,)).fetchone()

        for tag in tags:
            # get tag id. Assume it exists?
            tag_id, = db.execute('select rowid from tags where tag = ?',
                                 (tag,)).fetchone()

            c = db.execute('''delete from source_tag
            where source_id = ? and tag_id = ?''', (source_id, tag_id))

            db.commit()
            print(f'Deleted {c.rowcount} rows ({tag} from {source}')


@cli.command()
@click.argument('tags', nargs=-1)
def delete_tag(tags):
    """Delete each tag.

    This should also delete tags from sources by cascade.
    """
    for tag in tags:
        c = db.execute('delete from tags where tag = ?', (tag,))
        print(f'Deleted {c.rowcount} rows ({tag})')
    db.commit()


@cli.command()
@click.argument('tags', nargs=-1)
@click.option('-f', '--fmt', default='{{ source }}\n{{ extra["citation"] }}')
def show_tag(tags, fmt):
    """Show entries with tags.

    FMT is a jinja template for the output. You have variables of source, text
    and extra.

    I don't have good logic here, we just show all entries. I could probably get
    some basic and logic with sets, but mostly I assume for now you only want
    one tag, so this works. TODO: add something like boolean logic?

    """
    template = Template(fmt)
    for tag in tags:
        for row in db.execute('''select
        sources.source, sources.text, sources.extra
        from sources
        inner join source_tag on source_tag.source_id = sources.rowid
        inner join tags on source_tag.tag_id = tags.rowid
        where tags.tag = ?''', (tag,)).fetchall():
            source, text, extra = row
            extra = json.loads(extra)
            richprint(template.render(**locals()))


@cli.command()
def list_tags():
    """Print defined tags."""
    print('The following tags are defined.')
    for tag, in db.execute('select tag from tags').fetchall():
        print(tag)


##########
# Review #
##########

@cli.command()
@click.option('-s', '--since',  default='1 week ago')
@click.option('-f', '--fmt', default=None)
def review(since, fmt):
    """Review new entries added SINCE.

    SINCE should be something dateparser can handle.
    FMT is a jinja template for the output. Defaults to an org-mode template.
    """
    since = dateparser.parse(since).strftime("%Y-%m-%d")
    c = db.execute('''select source, text, extra from sources
    where date(date_added) > ?''', (since,)).fetchall()

    template = Template(fmt or '''* {{ source }}
:PROPERTIES:
:CITED_BY_COUNT: {{ data.get('cited_by_count', 0) }}
:END:

{{ text }}
        ''')

    for source, text, extra in c:
        data = json.loads(extra) or {}
        print(template.render(**locals()))


#############
# Searching #
#############

@cli.command()
def screenshot():
    """Do vector search from text in a screenshot.

    Use OCR to get text from an image on the clipboard (probably from a
    screenshot) and do a vector search on the text.
    """
    from PIL import ImageGrab
    import pytesseract
    os.environ["TOKENIZERS_PARALLELISM"] = "false"

    # Grab the image from the clipboard
    img = ImageGrab.grabclipboard()

    if img:
        text = pytesseract.image_to_string(img)
        print(f'Searching for "{text}"')
        vsearch(text)
    else:
        print("No image found in clipboard.")


def record():
    """Make a recording to a tempfile.

    Returns the audio filename it is recorded in.
    """
    import pyaudio
    import wave
    import threading

    # Parameters for recording
    FORMAT = pyaudio.paInt16
    CHANNELS = 2
    RATE = 44100
    CHUNK = 1024

    _, audio_file = tempfile.mkstemp(suffix='.wav')

    # Initialize PyAudio
    audio = pyaudio.PyAudio()

    # Function to record audio
    def record_audio():
        # Open stream
        stream = audio.open(format=FORMAT, channels=CHANNELS,
                            rate=RATE, input=True,
                            frames_per_buffer=CHUNK)
        print("Recording... Press Enter to stop.")

        frames = []

        # Record until Enter is pressed
        while not stop_recording.is_set():
            data = stream.read(CHUNK)
            frames.append(data)

        # Stop and close the stream
        stream.stop_stream()
        stream.close()

        # Save the recorded data as a WAV file
        wf = wave.open(audio_file, 'wb')
        wf.setnchannels(CHANNELS)
        wf.setsampwidth(audio.get_sample_size(FORMAT))
        wf.setframerate(RATE)
        wf.writeframes(b''.join(frames))
        wf.close()

    # Event to signal when to stop recording
    stop_recording = threading.Event()

    # Start recording in a separate thread
    recording_thread = threading.Thread(target=record_audio)
    recording_thread.start()

    # Wait for Enter key press to stop recording
    input()
    stop_recording.set()

    # Wait for the recording thread to finish
    recording_thread.join()

    # Terminate PyAudio
    audio.terminate()

    return audio_file


@cli.command()
@click.option('-p', '--playback', is_flag=True, help='Play audio back')
def audio(playback=False):
    """Record audio, convert it to text, and do a vector search on the text.

    The idea is nice, but the quality of transcription for scientific words is
    not that great. A better transcription library might make this more useful.
    """
    import speech_recognition as sr

    while True:

        afile = record()
        audio_file = sr.AudioFile(afile)
        r = sr.Recognizer()
        with audio_file as source:
            audio = r.listen(source)
            try:
                text = r.recognize_sphinx(audio)
                print('\n' + text + '\n')
            except sr.UnknownValueError:
                print("Could not understand audio.")

        if playback:
            import playsound
            playsound.playsound(afile, block=True)

        response = input('Is that what you want to search? ([y]/n/q): ')
        if response.lower().startswith('q'):
            return
        elif response.lower().startswith('n'):
            # record a new audio
            continue
        else:
            # move on to searching
            break

    vsearch(text)


@cli.command()
@click.argument('query', nargs=-1)
@click.option('-n',  default=3)
@click.option('-e', '--emacs', is_flag=True, default=False)
@click.option('-f', '--fmt', default=" {{ i }}. ({{ similarity|round(3) }}) {{ text }}\n\n")
@click.option('-x', '--cross-encode', is_flag=True, default=False)
def vsearch(query, n, emacs, fmt, cross_encode):
    """Do a vector search on QUERY.

    N is an integer for number of results to return
    EMACS is a flag for changing the output format for emacs

    FORMAT is a jinja template for the output. The variables you have access to
    are i, source, text, extra, similarity.

    CROSS_ENCODE is a boolean that resorts the results with a cross-encoder.

    """
    query = ' '.join(query)
    model = SentenceTransformer(config['embedding']['model'])
    emb = model.encode([query]).astype(np.float32).tobytes()

    c = db.execute('''select sources.source, sources.text,
    sources.extra, vector_distance_cos(?, embedding)
    from vector_top_k('embedding_idx', ?, ?)
    join sources on sources.rowid = id''', (emb, emb, n))

    results = c.fetchall()

    if cross_encode:
        import torch
        from sentence_transformers.cross_encoder import CrossEncoder

        # I don't know why I have to set the activation function here, but the
        # score is not 0..1 otherwise
        ce = CrossEncoder(config['embedding']['cross-encoder'],
                          default_activation_function=torch.nn.Sigmoid())
        scores = ce.predict([[query, text] for _, text, _, _ in results])
        # resort based on the scores
        results = [results[i] for i in np.argsort(scores)]

    if emacs:
        tmpl = ('( {% for source, text, extra, distance in results %}'
                '("({{ distance|round(3) }}) {{ text }}" . "{{ source }}") '
                '{% endfor %})')
        template = Template(tmpl)
        print(template.render(**locals()))
    else:
        for i, row in enumerate(results, 1):
            source, text, extra, similarity = row
            template = Template(fmt)
            richprint(template.render(**locals()))


@cli.command()
@click.argument('query', nargs=-1)
@click.option('-n',  default=3)
def fulltext(query, n):
    """Perform a fulltext search on litdb."""
    query = ' '.join(query)

    for source, text in db.execute('''select
    source, snippet(fulltext, 1, '', '', '', 16)
    from fulltext
    where text match ? order by rank limit ?''', (query, n)).fetchall():
        richprint(f"[link]{source}[/link]")
        richprint(text + '\n')


# Adapted from https://www.arsturn.com/blog/understanding-ollamas-embedding-models
@cli.command()
@click.argument('prompt', nargs=-1)
def gpt(prompt):
    """Run an ollama query with PROMPT."""
    t0 = time.time()
    prompt = ' '.join(prompt)
    model = SentenceTransformer(config['embedding']['model'])
    emb = model.encode([prompt]).astype(np.float32).tobytes()
    richprint(f'It took {time.time() - t0:1.1f} sec to embed the prompt')
    t0 = time.time()
    data = db.execute('''select sources.text
    from vector_top_k('embedding_idx', ?, 3)
    join sources on sources.rowid = id''', (emb,)).fetchall()

    richprint(f'It took  {time.time() - t0:1.1f} sec to get the top three docs')
    t0 = time.time()
    prompt = f"Using data: {data}. Respond to the prompt: {prompt}"
    output = ollama.generate(model="llama2",
                             prompt=prompt)
    richprint(output['response'])
    richprint(f'It took  {time.time() - t0:1.1f} '
              'sec to generate the response.')

    richprint('The text was generated using these references')
    for i, result in enumerate(data):
        richprint(f'{i:2d}. {result}\n')


@cli.command()
@click.argument('source')
@click.option('-f', '--fmt', default=None)
@click.option('-e', '--emacs', is_flag=True, default=False)
@click.option('-n',  default=3)
def similar(source, n, emacs, fmt):
    """Find N sources similar to SOURCE by vector similarity.

    if EMACS is truthy, the output is lisp for Emacs to read.
    FMT is a jinja template with access to source, text, and extra
    """
    emb, = db.execute('''select embedding from sources where source = ?''',
                      (source,)).fetchone()

    rows = db.execute('''select sources.source, sources.text, sources.extra
    from vector_top_k('embedding_idx', ?, ?)
    join sources on sources.rowid = id''', (emb, n + 1)).fetchall()[1:]
    # we do n + 1 because the first entry is always the source

    if emacs:
        template = Template('({% for source, text, extra in rows %}'
                            ' ("({{ text }}" . "{{ source }}")'
                            ' {% endfor %})')
        print(template.render(**locals()))
    else:
        template = Template(fmt or '{{ i }}. {{ source }}\n {{text}}\n\n')
        # print starting at index 1, the first item is always the source.
        for i, row in enumerate(rows, 1):
            source, text, extra = row
            richprint(template.render(**locals()))


@cli.command()
@click.argument('query', nargs=-1)
@click.option('-n', default=3)
@click.option('-m', '--max-steps', default=None)
def isearch(query, n=3, max_steps=None):
    """Perform an iterative search on QUERY.

    N is the number of results to return in each search.
    You will be prompted in each iteration if you want to continue or not.
    """
    query = ' '.join(query)
    model = SentenceTransformer(config['embedding']['model'])
    emb = model.encode([query]).astype(np.float32).tobytes()

    best = None

    steps = 0

    while True:

        results = db.execute('''select
        sources.source, sources.text,
        sources.extra, vector_distance_cos(?, embedding) as d
        from vector_top_k('embedding_idx', ?, ?)
        join sources on sources.rowid = id
        order by d''', (emb, emb, n)).fetchall()

        for source, text, extra, d in results:
            print(f'{d:1.3f}: {source}')

        steps += 1
        if steps == max_steps:
            break            

        current = [x[0] for x in results]  # sources

        # This means no change
        if current == best:
            print('Nothing new was found')
            break

        if input('Continue ([y]/n)?').lower().startswith('n'):
            break

        # something changed. add references and loop
        best = current
        for source in current:
            add_work(source, True, True, True)

    for source, text, extra, d in results:
        print(f'{d:1.3f}: {source}\n{text}')

    return results


@cli.command()
@click.argument('query', nargs=-1)
@click.option('-g', '--google', is_flag=True)
@click.option('-gs', '--google-scholar', is_flag=True)
@click.option('-pm', '--pubmed', is_flag=True)
@click.option('-ar', '--arxiv', is_flag=True)
@click.option('-cr', '--chemrxiv', is_flag=True)
@click.option('-br', '--biorxiv', is_flag=True)
@click.option('-a', '--all', is_flag=True)
def web(query, google, google_scholar, pubmed, arxiv, chemrxiv, biorxiv, all):
    """Open a web search for QUERY.

    We always do an OpenAlex search.

    If these are true, we also search here
    GOOGLE
    GOOGLE_SCHOLAR
    PUBMED
    ARXIV
    CHEMRXIV
    BIORXIV
    """
    query = ' '.join(query)

    if all:
        google, google_scholar = True, True, True,
        pubmed, arxiv, chemrxiv, biorxiv = True, True, True, True

    print(google, google_scholar, pubmed, arxiv)
    # This is to avoid some huggingface/tokenizer warning. I don't know why we
    # need to do it, but this code forks the process, and triggers that warning.
    os.environ['TOKENIZERS_PARALLELISM'] = 'false'

    oa = f'https://openalex.org/works?filter=default.search:{query}'
    webbrowser.open(oa)

    if google:
        url = f'https://www.google.com/search?q={query}'
        webbrowser.open(url)

    if google_scholar:
        url = f'https://scholar.google.com/scholar?q={query}'
        webbrowser.open(url)

    if pubmed:
        url = f'https://pubmed-ncbi-nlm-nih-gov.cmu.idm.oclc.org/?term={query}'
        webbrowser.open(url)

    if arxiv:
        url = f'https://arxiv.org/search/?query={query}'
        webbrowser.open(url)

    if chemrxiv:
        url = f'https://chemrxiv.org/engage/chemrxiv/search-dashboard?text={query}'
        webbrowser.open(url)

    if biorxiv:
        url = f'https://www.biorxiv.org/search/{query}'
        webbrowser.open(url)



###########
# Filters #
###########

@cli.command()
@click.argument('filter')
@click.option('-d', '--description')
def add_filter(filter, description=None):
    """Add an OpenAlex FILTER."""
    db.execute('insert into queries(filter, description) values (?, ?)',
               (filter, description))
    db.commit()


@cli.command()
@click.argument('filter')
def rm_filter(filter):
    """Remove an OpenAlex FILTER."""
    db.execute('delete from queries where filter = ?',
               (filter,))
    db.commit()


@cli.command()
def update_filters():
    """Update litdb using a filter with works from a created date."""
    filters = db.execute('''select filter, last_updated from queries''')
    for f, last_updated in filters.fetchall():
        update_filter(f, last_updated)


@cli.command()
def list_filters():
    """List the filters."""
    filters = db.execute('''select rowid, filter, description, last_updated
    from queries''')
    for rowid, f, description, last_updated in filters.fetchall():
        richprint(f'{rowid:3d}. {description or "None":30s} : {f} '
                  f'({last_updated})')


######################
# OpenAlex searching #
######################

@cli.command()
@click.argument('query', nargs=-1)
@click.option('-f', '--filter', '_filter', is_flag=True, default=False)
@click.option('-e', '--endpoint', default='works')
def openalex(query, _filter, endpoint):
    """Run an openalex query on FILTER.

    ENDPOINT should be one of works, authors, or another entity.

    This does not add anything to your database. It is to help you find starting
    points.

    To search text:
    litdb openalex "circular polymer"

    To find a journal id with a specific filter
    litdb openalex -e sources -f "display_name.search:Digital Discovery"

    """
    url = f'https://api.openalex.org/{endpoint}'

    query = ' '.join(query)
    if not _filter:
        query = f'default.search:{query}'

    next_cursor = '*'
    params = {'email': config['openalex']['email'],
              'api_key': config['openalex'].get('api_key'),
              'filter': query,
              'cursor': next_cursor}

    ok = False  # flag for large results
    while next_cursor:
        resp = requests.get(url, params)
        if resp.status_code != 200:
            print(resp.url)
            print(resp.text)
            return

        data = resp.json()

        count = data['meta']['count']
        if not ok and count > 100:
            resp = input(f'Found {count} results. Continue? ([y]/n)')
            if resp.lower().startswith('n'):
                return
            else:
                ok = True
        next_cursor = data['meta']['next_cursor']
        params.update(cursor=next_cursor)

        for result in data['results']:
            try:
                richprint(get_text(result))
                print()
            except:
                richprint(f'{result["id"]}: {result["display_name"]}')


########################################
# Convenience functions to add filters #
########################################

@cli.command()
@click.argument('name', nargs=-1)
def author_search(name):
    """Search OpenAlex for name.

    Uses the autocomplete endpoint to find an author's orcid.
    """
    auname = ' '.join(name)

    url = 'https://api.openalex.org/autocomplete/authors'

    from .openalex import get_data

    data = get_data(url,
                    params={'q': auname})

    for result in data['results']:
        richprint(f'- {result["display_name"]}\n  {result["hint"]} '
                  f'{result["external_id"]}\n\n')


@cli.command()
@click.argument('orcid')
@click.option('-r', '--remove', is_flag=True, help='remove')
def follow(orcid, remove=False):
    """Add a filter to follow orcid."""
    if not orcid.startswith('http'):
        orcid = f'https://orcid.org/{orcid}'

    # Seems like we should get their articles first.
    add_author(orcid)

    f = f'author.orcid:{orcid}'

    if remove:
        c = db.execute('''delete from queries where  filter = ?''',
                       (f,))
        db.commit()
        richprint(f'{c.rowcount} rows removed')
        return

    url = f'https://api.openalex.org/authors/{orcid}'
    data = get_data(url)
    name = data['display_name']

    today = datetime.date.today().strftime('%Y-%m-%d')
    db.execute('''insert or ignore into
    queries(filter, description, last_updated)
    values (?, ?, ?)''', (f, name, today))

    richprint(f'Following {name}: {orcid}')
    db.commit()


@cli.command()
@click.argument('query', nargs=-1)
@click.option('-r', '--remove', is_flag=True, help='remove')
def watch(query, remove=False):
    """Create a watch on query.

    QUERY: string, a filter for openalex.
    REMOVE: a flag to remove the query.
    """
    # First, we should make sure the filter is valid
    query = ' '.join(query)

    if remove:
        c = db.execute('''delete from queries where filter = ?''', (query,))
        db.commit()
        richprint(f'{c.rowcount} rows removed')
        return

    url = 'https://api.openalex.org/works'

    data = get_data(url, params={'filter': query})
    if len(data['results']) == 0:
        richprint(f"Sorry, {query} does not seem valid.")

    if remove:
        c = db.execute('''delete from queries where filter = ?''', (query,))
        richprint(f'Deleted {c.rowcount} rows')
        db.commit()
    else:
        c = db.execute('''insert or ignore into queries(filter, description)
        values (?, ?)''', (query,))
        richprint(f'Added {c.rowcount} rows')
        db.commit()
        richprint(f'Watching {query}')


@cli.command()
@click.argument('doi')
@click.option('-r', '--remove', is_flag=True, help='remove')
def citing(doi, remove=False):
    """Create a watch for articles that cite doi.

    REMOVE is a flag to remove the doi.
    """
    url = 'https://api.openalex.org/works'

    # We need an OpenAlex id
    f = f'doi:{doi}'

    data = get_data(url, params={'filter': f})
    if len(data['results']) == 0:
        richprint(f"Sorry, {doi} does not seem valid.")

    wid = data['results'][0]['id']

    if remove:
        c = db.execute('''delete from queries where filter = ?''',
                       (f'cites:{wid}',))
        db.commit()
        richprint(f'Deleted {c.rowcount} rows')
    else:
        c = db.execute('''insert or ignore into queries(filter, description)
        values (?, ?)''', (f'cites:{wid}', f'Citing papers for {doi}'))

        db.commit()
        richprint(f'Added {c.rowcount} rows')


@cli.command()
@click.argument('doi')
@click.option('-r', '--remove', is_flag=True, help='remove')
def related(doi, remove=False):
    """Create a watch for articles that are related to doi.

    REMOVE is a flag to remove the doi from queries.
    """
    url = 'https://api.openalex.org/works'

    # We need an OpenAlex id
    f = f'doi:{doi}'

    data = get_data(url, params={'filter': f})
    if len(data['results']) == 0:
        richprint(f"Sorry, {doi} does not seem valid.")

    wid = data['results'][0]['id']

    if remove:
        c = db.execute('''delete from queries where filter = ?''',
                       (f'related_to:{wid}',))
        db.commit()
        richprint(f'Deleted {c.rowcount} rows')
    else:
        c = db.execute('''insert or ignore into queries(filter, description)
        values (?, ?)''',
        (f'related_to:{wid}', f'Related papers for {doi}'))

        db.commit()
        richprint(f'Added {c.rowcount} rows')


#############
# Utilities #
#############

@cli.command()
@click.argument('sources', nargs=-1)
def bibtex(sources):
    """Generate bibtex entries for sources."""
    if not sources:
        sources = sys.stdin.read().strip().split()

    for source in sources:
        work, = db.execute('''select extra from sources where source = ?''',
                           (source,)).fetchone()
        richprint(dump_bibtex(json.loads(work)))


@cli.command()
@click.argument('sources', nargs=-1)
def citation(sources):
    """Generate citation strings for sources."""
    if not sources:
        sources = sys.stdin.read().strip().split()

    for i, source in enumerate(sources):
        citation, = db.execute('''select json_extract(extra, '$.citation')
        from sources where source = ?''', (source,)).fetchone()
        richprint(f'{i + 1:2d}. {citation}')


@cli.command()
@click.argument('doi')
def unpaywall(doi):
    """Use unpaywall to find PDFs for doi."""

    url = f'https://api.unpaywall.org/v2/{doi}'
    params = {'email': config['openalex']['email']}

    data = requests.get(url, params).json()
    richprint(f'{data["title"]}, {data.get("journal_name") or ""}')
    richprint(f'Is open access: {data.get("is_oa", False)}')

    for loc in data.get('oa_locations', []):
        richprint(loc.get('url_for_pdf') or loc.get('url_for_landing_page'))


@cli.command()
def about():
    """Summary statistics of your db."""
    dbf = root / config['database']['db']
    richprint(f'Your database is located at {dbf}')
    kb = 1024
    mb = 1024 * kb
    gb = 1024 * mb
    richprint(f'Database size: {os.path.getsize(dbf) / gb:1.2f} GB')
    nsources, = db.execute('select count(source) from sources').fetchone()
    richprint(f'You have {nsources} sources')


@cli.command()
@click.argument('sql')
def sql(sql):
    """Run the SQL command on the db."""
    for row in db.execute(sql).fetchall():
        richprint(row)


@cli.command()
@click.argument('sources', nargs=-1)
@click.option('-f', '--fmt', default='{{ source }}\n{{ text }}')
def show(sources, fmt):
    """Show the source.

    FMT is a jinja template with access to source, text, extra for each arg.
    """
    for src in sources:
        result = db.execute('''select source, text, extra from
        sources where source = ?''', (src,)).fetchone()

        if result:
            source, text, extra = result
            extra = json.loads(extra)
            richprint(Template(fmt).render(**locals()))
        else:
            print(f'Nothing found for {src}')


@cli.command()
@click.argument('source')
def open(source):
    """Open source."""
    if source.startswith('http'):
        webbrowser.open(source, new=2)
    elif source.endswith('.pdf'):
        webbrowser.open(f'file://{source}')
    else:
        webbrowser.open(f'file://{source}')

######################
# Academic functions #
######################

@cli.command()
@click.argument('orcid')
def coa(orcid):
    from .coa import get_coa
    get_coa(orcid)


@cli.command()
@click.argument('query', nargs=-1)
@click.option('-n', default=5)
def suggest_reviewers(query, n):
    """Suggest reviewers for QUERY.

    Use up to N similar documents. This is an iterative function, you will be
    prompted to expand the search.
    """
    query = ' '.join(query)
    model = SentenceTransformer(config['embedding']['model'])
    emb = model.encode([query]).astype(np.float32).tobytes()

    results = db.execute('''select
        sources.source, sources.text,
        sources.extra, vector_distance_cos(?, embedding) as d
        from vector_top_k('embedding_idx', ?, ?)
        join sources on sources.rowid = id
        order by d''', (emb, emb, n)).fetchall()

    for source, text, extra, d in results:
        richprint(f'{d:1.3f}: {source}')

    if input('Look for something better (n/y): ').lower().startswith('y'):
        results = isearch(query.split(' '), n=n)

    # Now collect the authors from the matching papers
    authors = []

    for i, row in enumerate(results):
        source, citation, extra, distance = row
                
        d = json.loads(extra)
    
        for authorship in d['authorships']:
            authors += [authorship['author']['id']]

    # get the unique ones    
    authors = set(authors)

    # Get author information
    data = []

    url = f'https://api.openalex.org/authors/'
    for batch in batched(authors, 50):
        url = f'https://api.openalex.org/authors?filter=id:{"|".join(batch)}'

        params = {'per-page': 50,
                  'email': config['openalex']['email']}

        r = get_data(url, params)

        for d in r['results']:
            lki = (d.get('last_known_institutions', [])[0].get('display_name') or
                   d['affiliations'][0]['institution']['display_name'])
            row = [d['display_name'], d['summary_stats']['h_index'], d['id'], lki]
            data += [row]
            
    # Sort and display the results
    data.sort(key = lambda row: row[1], reverse=True)

    print(tabulate.tabulate(data, headers=['name', 'h-index', 'oaid', 'institution'],
                            tablefmt='orgtbl'))

    print('\n' + 'From these papers:')
    for i, row in enumerate(results):
        source, citation, extra, distance = row
        print(f'{i + 1:2d}. {citation} (source)\n\n')
if __name__ == '__main__':
    cli()

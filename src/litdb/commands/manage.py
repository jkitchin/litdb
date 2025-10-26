"""Core database management commands.

Commands:
- init: Initialize litdb project
- add: Add sources to database
- remove: Remove sources
- index: Index file directories
- reindex: Reindex saved directories
- update_embeddings: Update embeddings for sources
"""

import os
import datetime
import pathlib

import bs4
import click
from docx import Document
import nbformat
from nbconvert import MarkdownExporter
from pptx import Presentation
import requests
from rich import print as richprint
from sentence_transformers import SentenceTransformer
from tqdm import tqdm

from ..utils import get_config, init_litdb
from ..db import get_db, add_source, add_work, add_author, add_bibtex
from ..pdf import add_pdf
from ..youtube import get_youtube_doc
from ..audio import is_audio_url, get_audio_text
from ..images import add_image, image_extensions


@click.command()
def init():
    """Initialize a litdb directory in the current working directory."""
    init_litdb()
    db = get_db()

    # Import about command to show statistics after init
    from ..cli import about

    with click.Context(about) as ctx:
        ctx.invoke(about)

    return db


@click.command()
@click.argument("sources", nargs=-1)
@click.option("--references", is_flag=True, help="Add references too.")
@click.option("--related", is_flag=True, help="Add related too.")
@click.option("--citing", is_flag=True, help="Add citing too.")
@click.option("--all", is_flag=True, help="Add references, related and citing.")
@click.option("-t", "--tag", "tags", multiple=True)
def add(
    sources,
    references=False,
    citing=False,
    related=False,
    all=False,
    verbose=False,
    tags=None,
):
    """Add WIDS to the db.

    REFERENCES, RELATED, CITING are flags to also add those for DOI sources. ALL
    is shorthand for all of those.

    SOURCES can be one or more of a doi or orcid, a pdf path, a url, bibtex
    file, or other kind of file assumed to be text.

    TAGS is a list of tags to add to the source.

    These are one time additions.

    """

    for source in tqdm(sources):
        # a work
        if source.startswith("10.") or "doi.org" in source:
            if source.startswith("10."):
                source = f"https://doi.org/{source}"

            if all:
                references, citing, related = True, True, True

            add_work(source, references, citing, related)

        # works from an author
        elif "orcid" in source or source.lower().startswith("https://openalex.org/a"):
            add_author(source)

        # a bibtex file
        elif source.endswith(".bib"):
            add_bibtex(source)

        # pdf
        elif source.endswith(".pdf"):
            source = os.path.abspath(source)
            add_pdf(source)

        # docx
        elif source.endswith(".docx"):
            source = os.path.abspath(source)
            doc = Document(source)
            add_source(source, "\n".join([para.text for para in doc.paragraphs]))

        # pptx
        elif source.endswith(".pptx"):
            source = os.path.abspath(source)
            prs = Presentation(source)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            add_source(source, "\n".join(text))

        # YouTube
        elif source.startswith("https") and "youtube" in source:
            text, citation = get_youtube_doc(source)
            add_source(source, text, {"citation": citation})

        # audio sources
        elif (source.startswith("http") and is_audio_url(source)) or source.endswith(
            ".mp3"
        ):
            add_source(source, get_audio_text(source))

        # local html
        elif not source.startswith("http") and source.endswith(".html"):
            source = os.path.abspath(source)
            with open(source) as f:
                text = f.read()
            soup = bs4.BeautifulSoup(text, features="lxml")
            add_source(source, soup.get_text())

        # a url
        elif source.startswith("http"):
            soup = bs4.BeautifulSoup(requests.get(source).text)
            add_source(source, soup.get_text())

        # ipynb
        elif source.endswith(".ipynb"):
            source = os.path.abspath(source)
            with open(source) as f:
                notebook = nbformat.read(f, as_version=4)

            # Create a Markdown exporter
            markdown_exporter = MarkdownExporter()

            # Convert the notebook to Markdown
            (body, resources) = markdown_exporter.from_notebook_node(notebook)

            add_source(source, body)

        # There are a lot of image extensions. I put this near the end so the
        # specific extensions are matched first.
        elif os.path.splitext(source)[1].lower() in image_extensions:
            add_image(source)

        # assume it is text
        else:
            source = os.path.abspath(source)
            with open(source) as f:
                text = f.read()
            add_source(source, text)

        # Handle tags if provided
        if tags:
            # TODO: Implement tagging during add
            pass


@click.command()
@click.argument("sources", nargs=-1)
def remove(sources):
    """Remove sources from litdb."""
    db = get_db()
    for source in sources:
        db.execute("delete from sources where source = ?", (source,))
        db.commit()


@click.command()
@click.argument("sources", nargs=-1)
def index(sources):
    """Index the directories in SOURCES.

    SOURCES is a list of directories.
    """
    db = get_db()
    for directory in sources:
        directory = pathlib.Path(directory).resolve()
        for fname in directory.rglob("*"):
            # for f in files:
            if fname.suffix in [
                ".pdf",
                ".docx",
                ".pptx",
                ".org",
                ".md",
                ".html",
                ".bib",
                ".ipynb",
            ]:
                fname = str(fname)

                # skip files we already have
                if db.execute(
                    """select source from sources where source = ?""", (fname,)
                ).fetchone():
                    continue

                with click.Context(add) as ctx:
                    print(fname)
                    ctx.invoke(add, sources=[fname])
                    print(f"Adding {fname}")

                    richprint(f"Added {fname}")

        last_updated = datetime.date.today().strftime("%Y-%m-%d")

        directory = str(directory)  # we need strings for the db
        if db.execute(
            """select path from directories where path = ?""", (directory,)
        ).fetchone():
            print(f"Updating {directory}")
            db.execute(
                """update directories set last_updated = ?
            where path = ?""",
                (last_updated, directory),
            )
        else:
            print(f"Inserting {directory}: {last_updated}")
            db.execute(
                """insert into directories(path, last_updated)
            values (?, ?)""",
                (directory, last_updated),
            )

        db.commit()


@click.command()
def reindex():
    """Reindex saved directories."""
    db = get_db()
    for (directory,) in db.execute("""select path from directories""").fetchall():
        print(f"Reindexing {directory}")
        index([directory])


@click.command(name="update-embeddings")
def update_embeddings():
    """Update all the embeddings in your db.

    The only reason you would do this is if you change the embedding model, or
    the way the chunks are sized in your config.

    """
    import numpy as np
    from langchain.text_splitter import RecursiveCharacterTextSplitter

    config = get_config()
    db = get_db()

    model = SentenceTransformer(config["embedding"]["model"])
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=config["embedding"]["chunk_size"],
        chunk_overlap=config["embedding"]["chunk_overlap"],
    )

    _, dim = model.encode(["test"]).shape

    # The point of this is to avoid deleting the database.
    db.execute("drop index if exists embedding_idx")
    db.execute("alter table sources drop embedding")
    db.execute(f"alter table sources add column embedding F32_BLOB({dim})")
    db.commit()

    for rowid, text in db.execute("select rowid, text from sources").fetchall():
        chunks = splitter.split_text(text)
        embedding = model.encode(chunks).mean(axis=0).astype(np.float32).tobytes()

        c = db.execute(
            "update sources set embedding = ? where rowid = ?", (embedding, rowid)
        )
        print(rowid, c.rowcount)

    # I don't know why this has to be here. I had it above, and no updates were
    # happening.
    db.execute(
        """create index if not exists embedding_idx ON sources (libsql_vector_idx(embedding))"""
    )
    db.commit()
